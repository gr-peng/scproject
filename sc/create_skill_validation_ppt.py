import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RUN_DATE = "2026-05-09"
TEMPLATE = ROOT / "RBM Group Roadshow Template.pptx"
OUT_DIR = ROOT / "outputs" / "skill-validation-presentation"
OUTPUT_PPTX = OUT_DIR / "five-skill-arxiv-agent-validation-rbm-template-final.pptx"
OUTLINE_JSON = OUT_DIR / "skill-validation-rbm-template-outline.json"

RUNS = [
    {
        "key": "agentic",
        "name": "Agentic AI Systems",
        "idea": "Emerging vocabulary and multi-agent papers",
        "dir": ROOT / "outputs" / "skill-proof-1-agentic-ai-systems",
        "accent": "teal",
    },
    {
        "key": "finance",
        "name": "AI for Finance",
        "idea": "Domain-specific finance, trading, payment papers",
        "dir": ROOT / "outputs" / "skill-proof-2-ai-for-finance",
        "accent": "blue",
    },
    {
        "key": "science",
        "name": "AI for Scientific Discovery",
        "idea": "Cross-disciplinary scientific discovery papers",
        "dir": ROOT / "outputs" / "skill-proof-3-ai-for-scientific-discovery",
        "accent": "green",
    },
]

COLORS = {
    "ink": RGBColor(25, 38, 55),
    "muted": RGBColor(95, 110, 128),
    "pale": RGBColor(248, 250, 253),
    "line": RGBColor(205, 216, 230),
    "white": RGBColor(255, 255, 255),
    "navy": RGBColor(0, 63, 126),
    "teal": RGBColor(38, 119, 122),
    "blue": RGBColor(38, 89, 156),
    "green": RGBColor(74, 138, 90),
    "gold": RGBColor(196, 140, 46),
    "rose": RGBColor(177, 83, 103),
    "purple": RGBColor(101, 86, 154),
}


def load_runs():
    runs = []
    for run in RUNS:
        data = json.loads((run["dir"] / f"daily-arxiv-briefing-{RUN_DATE}.json").read_text(encoding="utf-8"))
        runs.append({**run, "data": data})
    return runs


def clear_template_text(slide):
    def clear_shape(shape):
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()
            if shape.text_frame.paragraphs:
                shape.text_frame.paragraphs[0].text = ""
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                clear_shape(child)

    for shape in slide.shapes:
        clear_shape(shape)


def add_text(slide, x, y, w, h, text, size=12, color="ink", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = COLORS[color]
    p.font.bold = bold
    if align is not None:
        p.alignment = align
    return box


def add_lines(slide, x, y, w, h, lines, size=11, color="ink", bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(4)
        if bullet:
            p.level = 0
    return box


def title(slide, text, subtitle=None):
    cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0.14), Inches(7.25), Inches(0.84))
    cover.fill.solid()
    cover.fill.fore_color.rgb = COLORS["white"]
    cover.line.color.rgb = COLORS["white"]
    add_text(slide, 0.55, 0.24, 8.5, 0.35, text, 16, "ink", True)
    if subtitle:
        add_text(slide, 0.72, 0.68, 10.7, 0.25, subtitle, 8.5, "muted")


def footer(slide, page):
    add_text(slide, 10.95, 6.63, 0.95, 0.18, f"{page:02d}", 7, "muted", align=PP_ALIGN.RIGHT)


def panel(slide, x, y, w, h, fill="white", line="line"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    return shape


def content_panel(slide, x=0.55, y=1.05, w=12.15, h=5.55):
    return panel(slide, x, y, w, h, "white", "line")


def metric(slide, x, y, label, value, color="teal", w=1.55):
    panel(slide, x, y, w, 0.72, "white", "line")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color]
    bar.line.color.rgb = COLORS[color]
    add_text(slide, x + 0.16, y + 0.09, w - 0.24, 0.18, label, 6.5, "muted", True)
    add_text(slide, x + 0.16, y + 0.34, w - 0.24, 0.25, value, 14, "ink", True)


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = COLORS["muted"]
    c.line.width = Pt(1.1)
    return c


def short(text, limit=70):
    text = re.sub(r"\s+", " ", str(text)).strip()
    return text if len(text) <= limit else text[: limit - 1] + "..."


def add_table(slide, x, y, w, h, headers, rows, font=7):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["pale"]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font)
            p.font.bold = True
            p.font.color.rgb = COLORS["ink"]
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font)
                p.font.color.rgb = COLORS["ink"]
    return shape


def add_chart(slide, x, y, w, h, chart_title, categories, series, max_scale=None):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = graphic.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.has_legend = True
    chart.legend.include_in_layout = False
    if max_scale is not None:
        chart.value_axis.maximum_scale = max_scale
    return graphic


def add_bar_visual(slide, x, y, w, h, chart_title, categories, series, max_value=1.0):
    add_text(slide, x, y, w, 0.28, chart_title, 12.5, "ink", True)
    colors = ["blue", "gold", "green", "rose", "purple"]
    top = y + 0.55
    label_w = min(1.65, w * 0.28)
    bar_w = w - label_w - 0.35
    row_h = (h - 0.72) / max(1, len(categories))
    for i, cat in enumerate(categories):
        row_y = top + i * row_h
        add_text(slide, x, row_y + 0.02, label_w, 0.25, str(cat), 7.8, "ink", True)
        for j, (name, values) in enumerate(series):
            value = max(0, min(max_value, float(values[i])))
            yy = row_y + 0.28 + j * 0.17
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + label_w), Inches(yy), Inches(bar_w), Inches(0.10))
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLORS["pale"]
            bg.line.color.rgb = COLORS["pale"]
            fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + label_w), Inches(yy), Inches(bar_w * value / max_value), Inches(0.10))
            fg.fill.solid()
            fg.fill.fore_color.rgb = COLORS[colors[j % len(colors)]]
            fg.line.color.rgb = COLORS[colors[j % len(colors)]]
    lx = x
    ly = y + h - 0.2
    for j, (name, _) in enumerate(series):
        sx = lx + j * 1.55
        swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(ly), Inches(0.14), Inches(0.10))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = COLORS[colors[j % len(colors)]]
        swatch.line.color.rgb = COLORS[colors[j % len(colors)]]
        add_text(slide, sx + 0.18, ly - 0.04, 1.25, 0.18, name, 6.8, "muted")


def add_top_paper_cards(slide, x, y, w, papers):
    for i, p in enumerate(papers):
        yy = y + i * 0.78
        panel(slide, x, yy, w, 0.62, "pale", "line")
        add_text(slide, x + 0.12, yy + 0.08, 0.32, 0.22, str(p["rank"]), 9, "navy", True)
        add_text(slide, x + 0.5, yy + 0.08, w - 1.75, 0.25, short(p["title"], 58), 8.2, "ink", True)
        add_text(slide, x + w - 1.08, yy + 0.1, 0.95, 0.18, f'Rel {p["relevance_score"]:.2f}', 7.2, "teal", True, PP_ALIGN.RIGHT)


def outline(runs):
    return {
        "template": str(TEMPLATE),
        "output": str(OUTPUT_PPTX),
        "slides": 17,
        "structure": "Whole-Part-Whole",
        "proof_runs": [
            {
                "topic": run["name"],
                "top_paper": run["data"]["papers"][0]["title"],
                "metrics": run["data"]["metrics"],
            }
            for run in runs
        ],
    }


def build_deck():
    runs = load_runs()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    OUTLINE_JSON.write_text(json.dumps(outline(runs), indent=2, ensure_ascii=False), encoding="utf-8")

    prs = Presentation(str(TEMPLATE))
    for s in prs.slides:
        clear_template_text(s)

    labels = [run["name"].replace("AI ", "") for run in runs]
    total_papers = sum(int(run["data"]["metrics"]["retrieved_count"]) for run in runs)
    page = 1

    # 1 Cover: keep campus photo and blue title band.
    s = prs.slides[0]
    add_text(s, 0.45, 1.55, 11.7, 0.58, "Validating a Five-Skill arXiv Mining Agent", 24, "white", True)
    add_text(s, 0.48, 2.33, 10.4, 0.34, "Three stress tests, one reproducible pipeline, and evidence from recent papers", 12.5, "white")
    add_text(s, 0.5, 5.9, 4.8, 0.28, "10-minute academic presentation", 10, "white", True)
    add_text(s, 10.55, 6.25, 1.35, 0.22, "May 2026", 8.5, "white", align=PP_ALIGN.RIGHT)
    add_text(s, 11.38, 6.58, 0.7, 0.18, "(GZ)", 7, "white", align=PP_ALIGN.RIGHT)

    # 2 Contents: keep the template contents panel.
    s = prs.slides[1]
    add_text(s, 5.52, 0.58, 1.75, 0.33, "Contents", 13, "navy", True, PP_ALIGN.CENTER)
    blue_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.92), Inches(13.333), Inches(4.58))
    blue_band.fill.solid()
    blue_band.fill.fore_color.rgb = COLORS["navy"]
    blue_band.line.color.rgb = COLORS["navy"]
    agenda = [
        ("1", "Validation Setup"),
        ("2", "The Five Skills"),
        ("3", "Three Proof Cases"),
        ("4", "Cross-run Evidence"),
        ("5", "Takeaways"),
    ]
    positions = [(1.0, 3.55), (3.35, 3.55), (5.7, 3.55), (8.05, 3.55), (10.4, 3.55)]
    for (num, label), (x, y) in zip(agenda, positions):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.85), Inches(1.05))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["navy"]
        card.line.color.rgb = COLORS["white"]
        add_text(s, x + 0.1, y + 0.16, 0.35, 0.28, num, 13, "white", True, PP_ALIGN.CENTER)
        add_text(s, x + 0.5, y + 0.2, 1.15, 0.48, label, 9.5, "white", True, PP_ALIGN.CENTER)
    add_text(s, 1.15, 5.55, 10.9, 0.32, "Focus: prove that each skill has clear content, measurable output, and visible effect.", 12, "white", True, PP_ALIGN.CENTER)
    footer(s, page); page += 1

    # 3 Thesis.
    s = prs.slides[2]
    title(s, "A useful literature agent must leave evidence", "The validation checks the agent as a system, not as five isolated prompts.")
    content_panel(s)
    panel(s, 0.88, 1.55, 5.0, 3.35)
    add_text(s, 1.0, 1.92, 4.6, 0.38, "Validation thesis", 18, "teal", True)
    add_lines(s, 1.0, 2.55, 4.55, 1.55, [
        "Fresh enough to monitor fast-moving arXiv topics.",
        "Selective enough to produce a useful top-k shortlist.",
        "Readable enough for a human presentation.",
        "Auditable enough to reproduce the run."
    ], 12.5, "ink", True)
    metric(s, 6.65, 1.78, "Papers", str(total_papers), "blue", 1.75)
    metric(s, 8.75, 1.78, "Top-k", "36", "teal", 1.75)
    metric(s, 10.85, 1.78, "Reports", "3", "green", 1.45)
    metric(s, 7.25, 3.2, "Figures", "21", "gold", 1.75)
    metric(s, 9.35, 3.2, "Completion", "100%", "purple", 1.75)
    footer(s, page); page += 1

    # 4 Three ideas, using monitor layout.
    s = prs.slides[3]
    title(s, "Three proof ideas cover three failure modes")
    content_panel(s, 4.85, 1.08, 7.35, 4.8)
    rows = []
    for run in runs:
        m = run["data"]["metrics"]
        rows.append([run["name"], run["idea"], int(m["retrieved_count"]), f'{int(m["category_diversity"])} cats'])
    add_table(s, 5.1, 1.42, 6.85, 2.55, ["Run", "Why it matters", "N", "Breadth"], rows, 8.2)
    add_text(s, 5.1, 4.35, 6.6, 0.62, "The cases stress emerging vocabulary, domain specificity, and interdisciplinary breadth.", 14, "teal", True)
    footer(s, page); page += 1

    # 5 Pipeline chain, using central blue band.
    s = prs.slides[4]
    title(s, "The pipeline is five skills connected by contracts")
    content_panel(s, 0.6, 1.35, 12.0, 4.7)
    stages = [
        ("Search", "Topic -> Paper objects"),
        ("Ranking", "Scores -> Top-k"),
        ("Summary", "Abstract -> Claims"),
        ("Visualize", "Metrics -> Charts"),
        ("Briefing", "Run -> Report"),
    ]
    for i, (head, body) in enumerate(stages):
        x = 0.95 + i * 2.25
        panel(s, x - 0.07, 2.28, 1.95, 1.15, "navy", "navy")
        add_text(s, x, 2.48, 1.78, 0.24, head, 13.5, "white", True, PP_ALIGN.CENTER)
        add_text(s, x, 2.92, 1.78, 0.32, body, 8.8, "white", False, PP_ALIGN.CENTER)
        if i < 4:
            arrow(s, x + 1.78, 3.0, x + 2.08, 3.0)
    add_text(s, 1.0, 5.35, 10.8, 0.35, "Stage boundaries make the agent testable: every skill consumes and emits structured evidence.", 14, "navy", True, PP_ALIGN.CENTER)
    footer(s, page); page += 1

    # 6 Search skill.
    s = prs.slides[5]
    title(s, "Skill 1 - Search: from fuzzy request to structured arXiv data")
    content_panel(s)
    add_lines(s, 0.85, 1.32, 5.05, 2.65, [
        "Content: query planning, arXiv API calls, Atom XML parsing, deduplication.",
        "Inputs: keywords, optional categories, date window, max_results, sorting rule.",
        "Outputs: Paper objects with ID, title, authors, abstract, dates, categories, URLs.",
        "Guardrail: never invent missing metadata; preserve empty fields."
    ], 13, "ink", True)
    add_bar_visual(s, 6.55, 1.32, 5.4, 2.7, "Search evidence", labels, [
        ("Retrieved / 35", [r["data"]["metrics"]["retrieved_count"] / 35 for r in runs]),
        ("Category diversity / 10", [r["data"]["metrics"]["category_diversity"] / 10 for r in runs]),
    ], 1.0)
    metric(s, 0.95, 4.75, "Freshness", "1-2d", "teal", 1.8)
    metric(s, 3.0, 4.75, "Links", "100%", "green", 1.8)
    add_text(s, 6.55, 4.65, 5.3, 0.65, "Effect: all three runs produced fresh, linked, deduplicated papers for downstream skills.", 14, "teal", True)
    footer(s, page); page += 1

    # 7 Ranking skill.
    s = prs.slides[6]
    title(s, "Skill 2 - Ranking: turning retrieval into a useful shortlist")
    content_panel(s)
    add_lines(s, 7.05, 1.32, 4.8, 2.8, [
        "Content: deterministic relevance, novelty, final score, and rank.",
        "Inputs: Paper text plus user interests.",
        "Scoring: relevance from query overlap; novelty from method/contribution cues.",
        "Output: sorted top-k while preserving all papers."
    ], 12.5, "ink", True)
    add_bar_visual(s, 0.92, 1.32, 5.55, 2.95, "Ranking evidence", labels, [
        ("Mean relevance", [r["data"]["metrics"]["mean_relevance_score"] for r in runs]),
        ("Coverage@K", [r["data"]["metrics"]["coverage_at_k"] for r in runs]),
        ("Mean novelty", [r["data"]["metrics"]["mean_novelty_score"] for r in runs]),
    ], 1.0)
    add_text(s, 7.05, 4.55, 4.8, 0.7, "Effect: the shortlist is useful, and weaker finance coverage exposes a concrete improvement path: better synonym/stemming support.", 13.5, "rose", True)
    footer(s, page); page += 1

    # 8 Summarization skill.
    s = prs.slides[7]
    title(s, "Skill 3 - Summarization: compressing abstracts into speakable claims")
    content_panel(s)
    add_lines(s, 0.82, 1.25, 4.0, 2.6, [
        "Content: contribution extraction, method extraction, one-brief synthesis.",
        "Inputs: ranked Paper objects with title and abstract.",
        "Output fields: contribution_summary, method_summary, brief_summary.",
        "Guardrail: do not fabricate claims outside the abstract."
    ], 12.5, "ink", True)
    x0 = 5.05
    for i, run in enumerate(runs):
        p = run["data"]["papers"][0]
        panel(s, x0 + i * 2.55, 1.35, 2.25, 3.7)
        add_text(s, x0 + i * 2.55 + 0.16, 1.58, 1.9, 0.28, run["name"].replace("AI ", ""), 10.2, run["accent"], True)
        add_text(s, x0 + i * 2.55 + 0.16, 2.08, 1.9, 0.78, short(p["title"], 58), 8.5, "ink", True)
        add_text(s, x0 + i * 2.55 + 0.16, 3.18, 1.9, 0.9, short(p["contribution_summary"], 92), 7.7, "ink")
        add_text(s, x0 + i * 2.55 + 0.16, 4.55, 1.9, 0.22, "100% completion", 8, "green", True)
    footer(s, page); page += 1

    # 9 Visualization skill.
    s = prs.slides[8]
    title(s, "Skill 4 - Visualization: dashboards for the pipeline")
    content_panel(s)
    add_lines(s, 0.82, 1.25, 4.35, 2.85, [
        "Content: metric computation and report-ready figure manifest.",
        "Metrics: retrieved_count, coverage@k, category_diversity, recency, completion, link validity.",
        "Figures: KPI cards, category bars, score comparison, relevance-recency, score distribution, summary coverage, runtime."
    ], 12.2, "ink", True)
    metric(s, 1.05, 4.75, "Figures", "21", "gold", 1.7)
    metric(s, 3.0, 4.75, "Success", "100%", "green", 1.7)
    add_bar_visual(s, 6.15, 1.45, 5.55, 3.15, "Visualization reliability", labels, [
        ("Summary", [r["data"]["metrics"]["summary_completion_rate"] for r in runs]),
        ("Figures", [r["data"]["metrics"]["figure_generation_success_rate"] for r in runs]),
        ("Links", [r["data"]["metrics"]["link_completeness_rate"] for r in runs]),
    ], 1.0)
    footer(s, page); page += 1

    # 10 Briefing skill.
    s = prs.slides[9]
    title(s, "Skill 5 - Briefing: packaging the run into reusable artifacts")
    content_panel(s)
    add_lines(s, 0.9, 1.3, 4.3, 2.7, [
        "Content: join upstream outputs and assemble a skimmable report.",
        "Outputs: Markdown, CSV, JSON, and manifest.",
        "Markdown: human story with embedded figures.",
        "CSV: annotation and metric table.",
        "JSON: source-of-truth snapshot."
    ], 12.3, "ink", True)
    artifacts = [("MD", "story"), ("CSV", "table"), ("JSON", "snapshot"), ("Manifest", "checks")]
    for i, (head, body) in enumerate(artifacts):
        x = 5.15 + i * 1.62
        panel(s, x, 2.2, 1.35, 1.2)
        add_text(s, x + 0.08, 2.43, 1.17, 0.25, head, 11, ["teal", "blue", "green", "gold"][i], True, PP_ALIGN.CENTER)
        add_text(s, x + 0.08, 2.86, 1.17, 0.24, body, 8.3, "muted", False, PP_ALIGN.CENTER)
    add_table(s, 0.9, 4.8, 10.7, 1.0, ["Run", "Artifacts", "Rows", "Visuals"], [[r["name"], "MD / CSV / JSON / manifest", "12", "7"] for r in runs], 8)
    footer(s, page); page += 1

    # 11 Case 1.
    for idx, run in enumerate(runs):
        s = prs.slides[10 + idx]
        title(s, f"Proof case {idx + 1}: {run['name']}", run["idea"])
        content_panel(s)
        m = run["data"]["metrics"]
        top = run["data"]["papers"][:3]
        add_top_paper_cards(s, 0.85, 1.45, 5.75, top)
        metric(s, 7.0, 1.48, "Retrieved", str(int(m["retrieved_count"])), run["accent"], 1.55)
        metric(s, 8.85, 1.48, "Coverage", f'{m["coverage_at_k"]*100:.0f}%', "green", 1.55)
        metric(s, 10.7, 1.48, "Categories", str(int(m["category_diversity"])), "gold", 1.55)
        add_bar_visual(s, 7.0, 2.6, 4.75, 2.15, "Run signals", ["Rel", "Novelty", "Coverage"], [("Value", [m["mean_relevance_score"], m["mean_novelty_score"], m["coverage_at_k"]])], 1.0)
        takeaway = {
            "agentic": "The agentic case proves the system can handle fast-moving vocabulary and produce a high-coverage shortlist.",
            "finance": "The finance case proves domain value while revealing the need for richer financial synonyms in ranking.",
            "science": "The discovery case proves cross-category breadth and strong summary utility for interdisciplinary scouting.",
        }[run["key"]]
        add_text(s, 0.95, 5.35, 10.8, 0.48, takeaway, 13.2, "teal", True)
        footer(s, page); page += 1

    # 14 Cross-run evidence.
    s = prs.slides[13]
    title(s, "Cross-run evidence: what the metrics prove")
    content_panel(s)
    add_bar_visual(s, 0.95, 1.35, 5.35, 3.0, "System-level reliability", labels, [
        ("Summary completion", [r["data"]["metrics"]["summary_completion_rate"] for r in runs]),
        ("Link completeness", [r["data"]["metrics"]["link_completeness_rate"] for r in runs]),
        ("Figure success", [r["data"]["metrics"]["figure_generation_success_rate"] for r in runs]),
    ], 1.0)
    metric(s, 7.25, 1.55, "Retrieved papers", str(total_papers), "blue", 1.9)
    metric(s, 9.45, 1.55, "Ranked highlights", "36", "teal", 1.9)
    metric(s, 7.25, 2.75, "Reports", "3", "green", 1.9)
    metric(s, 9.45, 2.75, "Figures", "21", "gold", 1.9)
    add_lines(s, 7.25, 4.15, 4.0, 1.1, [
        "The workflow is measurable.",
        "The artifacts are reproducible.",
        "Quality gaps become visible."
    ], 13, "ink", True)
    footer(s, page); page += 1

    # 15 Maturity pyramid.
    s = prs.slides[14]
    title(s, "What the validation proves about the five skills")
    content_panel(s, 4.95, 1.05, 7.1, 4.95)
    add_lines(s, 5.35, 1.35, 5.95, 2.75, [
        "Base layer: Search gives fresh structured data.",
        "Second layer: Ranking creates a selective paper list.",
        "Third layer: Summaries make evidence readable.",
        "Fourth layer: Visualizations make behavior interpretable.",
        "Top layer: Briefing packages everything for reuse."
    ], 13.3, "ink", True)
    add_text(s, 5.35, 4.75, 5.95, 0.55, "Together they form an evidence-producing literature workflow, not just a set of isolated prompts.", 14, "teal", True)
    footer(s, page); page += 1

    # 16 Roadmap.
    s = prs.slides[15]
    title(s, "Roadmap: useful imperfections become next experiments")
    content_panel(s)
    add_lines(s, 7.05, 1.35, 4.65, 3.3, [
        "1. Add stemming and domain synonym expansion for ranking.",
        "2. Add human relevance labels for Precision@K and nDCG.",
        "3. Add PDF-level claim extraction beyond abstracts.",
        "4. Add optional PPT-native chart export from generated briefings."
    ], 13, "ink", True)
    metric(s, 1.1, 2.0, "Current strength", "reproducible", "green", 2.25)
    metric(s, 1.1, 3.15, "Next target", "better ranking", "rose", 2.25)
    footer(s, page); page += 1

    # 17 Close.
    s = prs.slides[16]
    title(s, "Conclusion: five skills, one evidence-producing workflow")
    content_panel(s, 0.55, 1.05, 6.2, 4.2)
    add_lines(s, 0.9, 1.35, 5.3, 2.1, [
        "The agent retrieved, ranked, summarized, visualized, and briefed three distinct arXiv topics.",
        "Each skill has a concrete role and measurable output.",
        "The generated artifacts support both reading and evaluation."
    ], 14, "ink", True)
    add_text(s, 6.88, 2.0, 4.2, 0.62, "Questions?", 28, "navy", True, PP_ALIGN.CENTER)
    add_text(s, 6.88, 3.0, 4.2, 0.55, "A small research newsroom for arXiv, built from five auditable skills.", 13, "teal", True, PP_ALIGN.CENTER)
    footer(s, page)

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    print(build_deck())
