import argparse
import json
from pathlib import Path


SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
DEFAULT_TEMPLATE = SKILL_DIR / "assets" / "academic-template.pptx"


def require_pptx():
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        raise SystemExit(
            "Missing dependency: python-pptx. Install it in the active Python environment "
            "before generating editable PPTX files."
        ) from exc
    return Presentation, CategoryChartData, XL_CHART_TYPE, MSO_SHAPE, PP_ALIGN, Inches, Pt


def text_frame(shape, paragraphs, font_size, color, bold=False):
    tf = shape.text_frame
    tf.clear()
    for index, text in enumerate(paragraphs):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.text = str(text)
        para.font.size = font_size
        para.font.color.rgb = color
        para.font.bold = bold
        para.space_after = 0
    return tf


def add_textbox(slide, x, y, w, h, text, font_size, color, bold=False):
    _, _, _, _, _, _, _, Inches, Pt = require_pptx()
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame(box, [text], Pt(font_size), color, bold=bold)
    return box


def add_bullets(slide, x, y, w, h, bullets, font_size, color):
    _, _, _, _, _, _, _, Inches, Pt = require_pptx()
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(bullets):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.text = str(item)
        para.level = 0
        para.font.size = Pt(font_size)
        para.font.color.rgb = color
        para.space_after = Pt(5)
    return box


def blank_slide(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def palette():
    from pptx.dml.color import RGBColor

    return {
        "ink": RGBColor(31, 41, 51),
        "muted": RGBColor(91, 111, 132),
        "accent": RGBColor(47, 111, 115),
        "blue": RGBColor(55, 103, 166),
        "gold": RGBColor(197, 139, 43),
        "panel": RGBColor(245, 247, 250),
        "white": RGBColor(255, 255, 255),
    }


def add_header(slide, title, subtitle=None):
    Presentation, _, _, _, _, _, _, Inches, Pt = require_pptx()
    colors = palette()
    add_textbox(slide, 0.62, 0.34, 11.0, 0.45, title, 22, colors["ink"], bold=True)
    if subtitle:
        add_textbox(slide, 0.64, 0.82, 10.5, 0.28, subtitle, 9, colors["muted"])


def add_title_slide(prs, spec):
    colors = palette()
    slide = blank_slide(prs)
    add_textbox(slide, 0.75, 1.65, 11.0, 1.0, spec.get("title", "Untitled Academic Talk"), 34, colors["ink"], bold=True)
    add_textbox(slide, 0.78, 2.78, 10.8, 0.5, spec.get("subtitle", ""), 17, colors["accent"])
    add_textbox(slide, 0.78, 3.52, 10.8, 0.4, spec.get("authors", ""), 12, colors["muted"])
    add_textbox(slide, 0.78, 6.65, 10.8, 0.28, spec.get("footer", ""), 9, colors["muted"])
    return slide


def add_section_slide(prs, spec):
    colors = palette()
    slide = blank_slide(prs)
    add_textbox(slide, 0.82, 2.2, 10.8, 0.9, spec.get("title", "Section"), 32, colors["accent"], bold=True)
    add_textbox(slide, 0.86, 3.16, 10.4, 0.6, spec.get("subtitle", spec.get("body", "")), 15, colors["ink"])
    return slide


def add_bullet_slide(prs, spec):
    colors = palette()
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Key Points"), spec.get("subtitle"))
    if spec.get("key_message"):
        add_textbox(slide, 0.74, 1.35, 10.9, 0.45, spec["key_message"], 15, colors["accent"], bold=True)
        y = 2.05
    else:
        y = 1.48
    add_bullets(slide, 0.94, y, 10.5, 4.3, spec.get("bullets", []), 16, colors["ink"])
    return slide


def add_two_column_slide(prs, spec):
    _, _, _, MSO_SHAPE, _, _, _, Inches, _ = require_pptx()
    colors = palette()
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Comparison"), spec.get("subtitle"))
    for index, side in enumerate(["left", "right"]):
        x = 0.72 + index * 5.9
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.45), Inches(5.25), Inches(4.7))
        box.fill.solid()
        box.fill.fore_color.rgb = colors["panel"]
        box.line.color.rgb = colors["muted"]
        add_textbox(slide, x + 0.25, 1.72, 4.7, 0.35, spec.get(f"{side}_title", side.title()), 16, colors["accent"], bold=True)
        add_bullets(slide, x + 0.35, 2.25, 4.55, 3.45, spec.get(f"{side}_bullets", []), 13, colors["ink"])
    return slide


def add_method_slide(prs, spec):
    _, _, _, MSO_SHAPE, _, _, _, Inches, _ = require_pptx()
    colors = palette()
    steps = spec.get("steps", [])
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Method Overview"), spec.get("subtitle"))
    count = max(1, len(steps))
    width = 10.5 / count
    for index, step in enumerate(steps):
        x = 0.78 + index * width
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(width - 0.18), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["panel"]
        shape.line.color.rgb = colors["accent"]
        add_textbox(slide, x + 0.18, 2.42, width - 0.55, 0.4, f"{index + 1}. {step}", 12, colors["ink"], bold=True)
        if index < count - 1:
            add_textbox(slide, x + width - 0.18, 2.54, 0.35, 0.3, "->", 16, colors["accent"], bold=True)
    add_bullets(slide, 0.9, 4.15, 10.4, 1.6, spec.get("bullets", []), 14, colors["ink"])
    return slide


def add_table_slide(prs, spec):
    _, _, _, _, _, _, _, Inches, Pt = require_pptx()
    colors = palette()
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Table"), spec.get("subtitle"))
    rows = spec.get("rows", [])
    headers = spec.get("headers", [])
    table_rows = max(1, len(rows) + (1 if headers else 0))
    table_cols = max(1, len(headers) or max([len(r) for r in rows], default=1))
    shape = slide.shapes.add_table(table_rows, table_cols, Inches(0.75), Inches(1.55), Inches(11.0), Inches(4.6))
    table = shape.table
    all_rows = ([headers] if headers else []) + rows
    for r, row in enumerate(all_rows):
        for c in range(table_cols):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10 if r else 11)
                para.font.bold = r == 0
                para.font.color.rgb = colors["ink"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["panel"] if r == 0 else colors["white"]
    return slide


def add_results_slide(prs, spec):
    Presentation, CategoryChartData, XL_CHART_TYPE, _, _, _, _, Inches, _ = require_pptx()
    colors = palette()
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Results"), spec.get("key_message"))
    chart = spec.get("chart")
    if chart:
        chart_data = CategoryChartData()
        chart_data.categories = chart.get("categories", [])
        for series in chart.get("series", []):
            chart_data.add_series(series.get("name", "Series"), series.get("values", []))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8),
            Inches(1.65),
            Inches(6.3),
            Inches(3.8),
            chart_data,
        )
        add_bullets(slide, 7.45, 1.72, 4.0, 3.8, spec.get("bullets", []), 13, colors["ink"])
    else:
        add_bullets(slide, 0.9, 1.6, 10.4, 4.2, spec.get("bullets", []), 15, colors["ink"])
    return slide


def add_takeaways_slide(prs, spec):
    _, _, _, MSO_SHAPE, _, _, _, Inches, _ = require_pptx()
    colors = palette()
    slide = blank_slide(prs)
    add_header(slide, spec.get("title", "Takeaways"), spec.get("subtitle"))
    for index, item in enumerate(spec.get("bullets", [])[:4]):
        y = 1.45 + index * 1.12
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(y), Inches(10.4), Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["panel"]
        shape.line.color.rgb = colors["accent"]
        add_textbox(slide, 1.18, y + 0.18, 9.6, 0.32, item, 15, colors["ink"], bold=True)
    return slide


def add_qa_slide(prs, spec):
    colors = palette()
    slide = blank_slide(prs)
    add_textbox(slide, 0.9, 2.35, 10.4, 0.8, spec.get("title", "Questions?"), 38, colors["accent"], bold=True)
    add_textbox(slide, 0.94, 3.28, 10.2, 0.4, spec.get("subtitle", spec.get("contact", "")), 14, colors["muted"])
    return slide


def build_deck(outline_path, output_path, template_path=DEFAULT_TEMPLATE):
    Presentation, *_ = require_pptx()
    with open(outline_path, "r", encoding="utf-8") as handle:
        outline = json.load(handle)
    prs = Presentation(str(template_path)) if Path(template_path).exists() else Presentation()
    slides = outline.get("slides", [])
    if not slides:
        slides = [{"layout": "title", "title": outline.get("title", "Untitled Academic Talk")}]
    handlers = {
        "title": add_title_slide,
        "section": add_section_slide,
        "bullets": add_bullet_slide,
        "two_column": add_two_column_slide,
        "method": add_method_slide,
        "results": add_results_slide,
        "table": add_table_slide,
        "takeaways": add_takeaways_slide,
        "qa": add_qa_slide,
    }
    for slide_spec in slides:
        handler = handlers.get(slide_spec.get("layout", "bullets"), add_bullet_slide)
        handler(prs, slide_spec)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Build an editable academic PPTX deck from a JSON outline.")
    parser.add_argument("outline_json", help="Path to JSON outline.")
    parser.add_argument("output_pptx", help="Path for generated PPTX.")
    parser.add_argument("--template", default=str(DEFAULT_TEMPLATE), help="PPTX template path.")
    args = parser.parse_args()
    build_deck(args.outline_json, args.output_pptx, args.template)


if __name__ == "__main__":
    main()
